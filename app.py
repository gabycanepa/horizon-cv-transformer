import streamlit as st
import os
import json
import time
import fitz  # PyMuPDF
from google import genai
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import tempfile
from io import BytesIO

# ==========================================
# CONFIGURACIÓN DE PÁGINA
# ==========================================
st.set_page_config(
    page_title="Transformador Horizon CV",
    page_icon="🎨",
    layout="centered"
)

st.markdown("""
<style>
    .main-header {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 2rem;
        border-radius: 10px;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        background: #d4edda;
        border: 2px solid #28a745;
        padding: 1.5rem;
        border-radius: 10px;
        margin-top: 1rem;
    }
    .stButton>button {
        width: 100%;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        font-size: 18px;
        font-weight: bold;
        padding: 0.75rem;
        border-radius: 10px;
        border: none;
    }
    .stButton>button:hover {
        background: linear-gradient(135deg, #764ba2 0%, #667eea 100%);
    }
</style>
""", unsafe_allow_html=True)

if "reset_key" not in st.session_state:
    st.session_state.reset_key = 0

def reset_app():
    st.session_state.reset_key += 1

# ==========================================
# CONFIGURACIÓN DE GEMINI
# ==========================================
API_KEY = st.secrets["GEMINI_API_KEY"]
client = genai.Client(api_key=API_KEY)
MODELO = "gemini-2.0-flash"

# ==========================================
# HELPERS
# ==========================================
def safe_str(value, default=""):
    if value is None:
        return default
    if isinstance(value, str):
        return value.strip()
    return str(value).strip()

def safe_list(value):
    return value if isinstance(value, list) else []

def normalizar_datos(datos):
    if isinstance(datos, list):
        datos = datos[0] if datos else {}

    if not isinstance(datos, dict):
        datos = {}

    contacto = datos.get("contacto")
    if not isinstance(contacto, dict):
        contacto = {
            "telefono": "",
            "email": "",
            "ubicacion": "",
            "linkedin": ""
        }

    experiencias = safe_list(datos.get("experiencias"))
    experiencias_normalizadas = []

    for exp in experiencias:
        if not isinstance(exp, dict):
            continue
        experiencias_normalizadas.append({
            "empresa": safe_str(exp.get("empresa")),
            "puesto": safe_str(exp.get("puesto")),
            "periodo": safe_str(exp.get("periodo")),
            "descripcion": safe_str(exp.get("descripcion"))
        })

    return {
        "nombre": safe_str(datos.get("nombre"), "N/A"),
        "rol": safe_str(datos.get("rol")),
        "contacto": {
            "telefono": safe_str(contacto.get("telefono")),
            "email": safe_str(contacto.get("email")),
            "ubicacion": safe_str(contacto.get("ubicacion")),
            "linkedin": safe_str(contacto.get("linkedin"))
        },
        "perfil": safe_str(datos.get("perfil")),
        "experiencias": experiencias_normalizadas,
        "educacion": safe_str(datos.get("educacion")),
        "habilidades_tecnicas": safe_str(datos.get("habilidades_tecnicas")),
        "herramientas": safe_str(datos.get("herramientas")),
        "expertise": safe_str(datos.get("expertise")),
        "certificaciones": safe_str(datos.get("certificaciones")),
        "idiomas": safe_str(datos.get("idiomas"))
    }

# ==========================================
# FUNCIONES DE PROCESAMIENTO
# ==========================================
def extraer_foto_y_texto(pdf_bytes):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    texto_completo = ""
    foto_bytes = None

    # 1. Extraer texto completo
    for page in doc:
        texto_completo += page.get_text()

    # 2. Intentar extraer imagen incrustada
    for page in doc:
        images = page.get_images(full=True)
        if images:
            mejor = max(
                images,
                key=lambda img: doc.extract_image(img[0])["width"] * doc.extract_image(img[0])["height"]
            )
            base_image = doc.extract_image(mejor[0])
            if base_image["width"] > 50 and base_image["height"] > 50:
                foto_bytes = base_image["image"]
                break

    # 3. Fallback: renderizar esquina superior izquierda
    if not foto_bytes and len(doc) > 0:
        page = doc[0]
        page_rect = page.rect
        clip = fitz.Rect(0, 0, page_rect.width * 0.22, page_rect.height * 0.30)
        mat = fitz.Matrix(3, 3)
        pix = page.get_pixmap(matrix=mat, clip=clip)
        foto_bytes = pix.tobytes("png")

    doc.close()
    return texto_completo, foto_bytes

def redactar_con_gemini(texto_cv):
    prompt = f"""
    Eres un transcriptor de datos de alta fidelidad para Horizon Consulting.
    Tu única misión es NO PERDER NINGUNA EXPERIENCIA LABORAL.

    Devuelve EXCLUSIVAMENTE JSON válido con esta estructura:
    {{
      "nombre": "",
      "rol": "",
      "contacto": {{
        "telefono": "",
        "email": "",
        "ubicacion": "",
        "linkedin": ""
      }},
      "perfil": "",
      "experiencias": [
        {{
          "empresa": "",
          "puesto": "",
          "periodo": "",
          "descripcion": ""
        }}
      ],
      "educacion": "",
      "habilidades_tecnicas": "",
      "herramientas": "",
      "expertise": "",
      "certificaciones": "",
      "idiomas": ""
    }}

    Reglas:
    - Nunca uses null.
    - Si falta un dato, devuelve "".
    - "experiencias" siempre debe ser una lista.
    - "descripcion" siempre debe ser string.
    - No agregues texto fuera del JSON.

    CV A PROCESAR:
    {texto_cv}
    """

    for intento in range(3):
        try:
            response = client.models.generate_content(
                model=MODELO,
                contents=prompt,
                config={"response_mime_type": "application/json"}
            )
            resultado = json.loads(response.text)
            return normalizar_datos(resultado)
        except Exception as e:
            if "429" in str(e) and intento < 2:
                time.sleep(5)
                continue
            raise e

def ajustar_fuente(shape, size=9):
    if hasattr(shape, "text_frame"):
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(size)

def llenar_shape_con_titulos(slide, placeholder, texto, title_size=Pt(11), body_size=Pt(9.5), font_color=(0, 0, 0)):
    texto = safe_str(texto)

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue

        try:
            current_text = safe_str(shape.text)
            if placeholder not in current_text:
                continue
        except Exception:
            continue

        tf = shape.text_frame
        tf.text = ""

        lines = [ln.rstrip() for ln in texto.splitlines() if ln.strip() != ""]

        if not lines:
            return True

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line

            stripped = line.strip()
            is_title = (
                stripped.endswith(":")
                or (stripped.isupper() and len(stripped) > 2)
            )

            if not p.runs:
                p.add_run()

            for run in p.runs:
                run.font.bold = is_title
                run.font.size = title_size if is_title else body_size
                run.font.color.rgb = RGBColor(*font_color)

        return True

    return False

def llenar_experiencias(slide, placeholder, experiencias, font_color=(0, 0, 0), titulo="EXPERIENCIA LABORAL:"):
    experiencias = safe_list(experiencias)

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue

        try:
            current_text = safe_str(shape.text)
            if placeholder not in current_text:
                continue
        except Exception:
            continue

        tf = shape.text_frame
        tf.text = ""

        p_titulo = tf.paragraphs[0]
        p_titulo.text = titulo
        if not p_titulo.runs:
            p_titulo.add_run()

        for run in p_titulo.runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(*font_color)

        for exp in experiencias:
            if not isinstance(exp, dict):
                continue

            empresa = safe_str(exp.get("empresa"), "?")
            puesto = safe_str(exp.get("puesto"), "?")
            periodo = safe_str(exp.get("periodo"), "?")
            desc = safe_str(exp.get("descripcion"), "")

            tf.add_paragraph().text = ""

            p_h = tf.add_paragraph()
            p_h.text = f"• {empresa} | {puesto} ({periodo})"
            if not p_h.runs:
                p_h.add_run()

            for run in p_h.runs:
                run.font.bold = True
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(*font_color)

            if desc:
                for d_line in desc.splitlines():
                    if d_line.strip():
                        p_d = tf.add_paragraph()
                        p_d.text = f"  {d_line.strip()}"
                        if not p_d.runs:
                            p_d.add_run()
                        for run in p_d.runs:
                            run.font.size = Pt(9.5)
                            run.font.color.rgb = RGBColor(*font_color)

        return True

    return False

def eliminar_cuadro_foto(slide):
    shapes_to_delete = []
    for s in slide.shapes:
        try:
            if hasattr(s, "text") and "FOTO" in safe_str(s.text).upper():
                shapes_to_delete.append(s)
        except Exception:
            pass

    for shape in shapes_to_delete:
        sp = shape.element
        sp.getparent().remove(sp)

def agregar_foto(slide, foto_bytes):
    if foto_bytes:
        eliminar_cuadro_foto(slide)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp.write(foto_bytes)
            tmp_path = tmp.name

        slide.shapes.add_picture(tmp_path, Inches(0.2), Inches(0.3), width=Inches(1.5))
        os.unlink(tmp_path)

def actualizar_encabezado(slide, nombre, rol):
    nombre = safe_str(nombre, "N/A")
    rol = safe_str(rol)

    for shape in slide.shapes:
        try:
            if not hasattr(shape, "text"):
                continue

            texto_shape = safe_str(shape.text)

            if any(x in texto_shape for x in ["{{NOMBRE}}", "COLABORADOR PROPUESTO", "Colaborador propuesto", "NOMBRE"]):
                tf = shape.text_frame
                tf.text = ""

                p1 = tf.paragraphs[0]
                p1.text = f"COLABORADOR PROPUESTO – {nombre}"
                if not p1.runs:
                    p1.add_run()
                for run in p1.runs:
                    run.font.bold = True
                    run.font.size = Pt(13)
                    run.font.color.rgb = RGBColor(88, 24, 139)

                p2 = tf.add_paragraph()
                p2.text = rol
                if not p2.runs:
                    p2.add_run()
                for run in p2.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(88, 24, 139)

                return
        except Exception:
            continue

def generar_pptx(datos, template_bytes, foto_bytes):
    datos = normalizar_datos(datos)

    prs = Presentation(BytesIO(template_bytes))
    nombre = safe_str(datos.get("nombre"), "N/A")
    rol = safe_str(datos.get("rol"))
    exps = safe_list(datos.get("experiencias"))

    if len(prs.slides) == 0:
        raise ValueError("El template PPTX no contiene slides.")

    # ── SLIDE 1 ──────────────────────────────────────────────
    s1 = prs.slides[0]
    actualizar_encabezado(s1, nombre, rol)

    llenar_shape_con_titulos(s1, "{{PUESTO}}", rol, Pt(11), Pt(10), font_color=(88, 24, 139))

    llenar_shape_con_titulos(
        s1,
        "{{PERFIL}}",
        f"PERFIL:\n{safe_str(datos.get('perfil'))}",
        Pt(11),
        Pt(9.5),
        font_color=(255, 255, 255)
    )

    edu = f"EDUCACIÓN:\n{safe_str(datos.get('educacion'))}"
    if safe_str(datos.get("certificaciones")):
        edu += f"\n\nCERTIFICACIONES:\n{safe_str(datos.get('certificaciones'))}"

    llenar_shape_con_titulos(
        s1,
        "{{EDUCACIÓN}}",
        edu,
        Pt(11),
        Pt(9.5),
        font_color=(255, 255, 255)
    )

    hab = f"HABILIDADES:\n{safe_str(datos.get('habilidades_tecnicas'))}"
    if safe_str(datos.get("herramientas")):
        hab += f"\n\nHERRAMIENTAS:\n{safe_str(datos.get('herramientas'))}"
    if safe_str(datos.get("expertise")):
        hab += f"\n\nEXPERTISE:\n{safe_str(datos.get('expertise'))}"

    llenar_shape_con_titulos(
        s1,
        "{{HABILIDADES}}",
        hab,
        Pt(11),
        Pt(9.5),
        font_color=(255, 255, 255)
    )

    llenar_shape_con_titulos(
        s1,
        "{{IDIOMAS}}",
        f"IDIOMAS:\n{safe_str(datos.get('idiomas'))}",
        Pt(11),
        Pt(9.5),
        font_color=(255, 255, 255)
    )

    if not llenar_experiencias(s1, "{{EXPERIENCIA_1_2}}", exps[:2], font_color=(0, 0, 0)):
        llenar_experiencias(s1, "{{EXPERIENCIA_RESTO}}", exps[:2], font_color=(0, 0, 0))

    agregar_foto(s1, foto_bytes)

    # ── SLIDES 2, 3, ... ─────────────────────────────────────
    rem = exps[2:]
    for i in range(1, len(prs.slides)):
        slide = prs.slides[i]
        actualizar_encabezado(slide, nombre, rol)
        llenar_shape_con_titulos(slide, "{{PUESTO}}", rol, Pt(11), Pt(10), font_color=(88, 24, 139))

        if rem:
            llenar_experiencias(
                slide,
                "{{EXPERIENCIA_RESTO}}",
                rem[:4],
                font_color=(0, 0, 0),
                titulo="EXPERIENCIA LABORAL (Cont.):"
            )
            rem = rem[4:]

        agregar_foto(slide, foto_bytes)

    out = BytesIO()
    prs.save(out)
    out.seek(0)

    return out, nombre, len(exps)

# ==========================================
# INTERFAZ DE STREAMLIT
# ==========================================
st.markdown(
    '<div class="main-header"><h1>🎨 TRANSFORMADOR HORIZON CV</h1><p>Convierte CVs al formato Horizon automáticamente</p></div>',
    unsafe_allow_html=True
)

with st.expander("ℹ️ Instrucciones de uso"):
    st.markdown("1. Sube el CV del candidato (.pdf)\n2. Haz clic en 'Transformar'\n3. Descarga el resultado.")

st.markdown("---")

col_main, col_btn = st.columns([8, 1])
with col_main:
    st.markdown("### 📁 Paso 1: Template Horizon")
with col_btn:
    st.button("🧹", on_click=reset_app, key="btn_limpiar")

DEFAULT_TEMPLATE = "CV HORIZON-MODELO 3.pptx"
template_bytes = None

template_file = st.file_uploader(
    "Sube un nuevo .pptx si quieres cambiar el modelo",
    type=["pptx"],
    key=f"t_{st.session_state.reset_key}"
)

if template_file:
    template_bytes = template_file.read()
    st.success(f"✓ Usando template subido: {template_file.name}")
elif os.path.exists(DEFAULT_TEMPLATE):
    with open(DEFAULT_TEMPLATE, "rb") as f:
        template_bytes = f.read()
    st.info(f"ℹ️ Usando modelo estándar: {DEFAULT_TEMPLATE}")
else:
    st.warning("⚠️ Sube un template .pptx para comenzar.")

st.markdown("---")

st.markdown("### 📄 Paso 2: Subir CV del Candidato")
cv_file = st.file_uploader(
    "Selecciona el archivo .pdf del CV",
    type=["pdf"],
    key=f"cv_{st.session_state.reset_key}"
)

if cv_file:
    st.success(f"✓ CV cargado: {cv_file.name}")

st.markdown("---")

if template_bytes and cv_file:
    if st.button("🚀 TRANSFORMAR A FORMATO HORIZON"):
        with st.spinner("⏳ Procesando..."):
            try:
                cv_bytes = cv_file.read()

                st.info("✓ Extrayendo información del PDF...")
                texto, foto_bytes = extraer_foto_y_texto(cv_bytes)

                st.info("✓ Analizando con Gemini...")
                datos_json = redactar_con_gemini(texto)

                st.info("✓ Generando presentación Horizon...")
                output_pptx, nombre, num_exp = generar_pptx(datos_json, template_bytes, foto_bytes)

                nombre_archivo = safe_str(nombre, "Candidato").replace(" ", "_")

                st.markdown(f"""
                <div class="success-box">
                    <h3 style='color: #155724; margin: 0;'>✅ ¡TRANSFORMACIÓN COMPLETADA!</h3>
                    <p style='margin: 10px 0;'><strong>Candidato:</strong> {safe_str(nombre, "N/A")}</p>
                    <p style='margin: 10px 0;'><strong>Experiencias detectadas:</strong> {num_exp}</p>
                </div>
                """, unsafe_allow_html=True)

                st.download_button(
                    label="📥 DESCARGAR CV HORIZON",
                    data=output_pptx,
                    file_name=f"CV_Horizon_{nombre_archivo}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            except Exception as e:
                st.error(f"❌ Error durante el procesamiento: {str(e)}")
                st.exception(e)
else:
    st.info("👆 Por favor sube el CV para continuar")

st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; padding: 1rem;'>
    <p>Desarrollado por Horizon Consulting | Powered by Gemini 2.0 Flash</p>
</div>
""", unsafe_allow_html=True)
